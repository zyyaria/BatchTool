# Copyright (C) 2026 张小鱼
# SPDX-License-Identifier: AGPL-3.0-or-later

import os
import io
import sys
import shutil
import subprocess
import urllib.parse
from PySide6.QtCore import Signal
from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel, QComboBox, QSpinBox,
    QLineEdit, QCheckBox, QSizePolicy, QButtonGroup, QPushButton
)
from core.utils import parse_page_range

try:
    import fitz
except ImportError:
    fitz = None
try:
    from PIL import Image
except ImportError:
    Image = None
try:
    from pdf2docx import Converter
except ImportError:
    Converter = None

if sys.platform.startswith("win"):
    try:
        import comtypes.client
        import pythoncom
        wdFormatPDF = 17
        wdFormatDoc = 0
        wdFormatDocx = 16
        ppSaveAsPDF = 32
        ppSaveAsPPT = 1
        ppSaveAsPPTX = 24
        xlTypePDF = 0
        xlExcel8 = 56
        xlOpenXMLWorkbook = 51
    except ImportError:
        comtypes = None
        pythoncom = None
else:
    comtypes = None
    pythoncom = None


class ConvertPanel(QWidget):
    changed = Signal()

    def __init__(self):
        """初始化设置面板"""
        super().__init__()
        layout = QVBoxLayout(self)
        layout.setSpacing(8)

        row_mode = QHBoxLayout()
        self.mode_group = QButtonGroup(self)
        self.btn_pdf = QPushButton("PDF 转换")
        self.btn_pdf.setCheckable(True)
        self.btn_pdf.setChecked(True)
        self.btn_pdf.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        self.btn_office = QPushButton("Office 互转")
        self.btn_office.setCheckable(True)
        self.btn_office.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        self.mode_group.addButton(self.btn_pdf)
        self.mode_group.addButton(self.btn_office)
        row_mode.addWidget(QLabel("模式:"))
        row_mode.addWidget(self.btn_pdf, 1)
        row_mode.addWidget(self.btn_office, 1)
        layout.addLayout(row_mode)

        row_format = QHBoxLayout()
        self.format_combo = QComboBox()
        self.format_combo.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        row_format.addWidget(QLabel("目标格式:"))
        row_format.addWidget(self.format_combo, 1)
        layout.addLayout(row_format)

        self.range_widget = QWidget()
        row_range = QHBoxLayout(self.range_widget)
        row_range.setContentsMargins(0, 0, 0, 0)
        self.range_edit = QLineEdit()
        self.range_edit.setPlaceholderText("1-3,5,8-（留空=全部）")
        self.range_edit.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        self.dpi_label = QLabel("DPI:")
        self.dpi_spin = QSpinBox()
        self.dpi_spin.setRange(72, 600)
        self.dpi_spin.setValue(72)
        self.dpi_spin.setSuffix(" ppi")
        self.dpi_spin.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Fixed)
        self.xlsx_check = QCheckBox("合并所有表")
        self.xlsx_check.setChecked(False)
        self.xlsx_check.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Fixed)
        row_range.addWidget(QLabel("页面范围:"))
        row_range.addWidget(self.range_edit, 1)
        row_range.addWidget(self.dpi_label)
        row_range.addWidget(self.dpi_spin, 1)
        row_range.addWidget(self.xlsx_check)
        layout.addWidget(self.range_widget)

        layout.addStretch()

        self.btn_pdf.toggled.connect(self._on_mode_changed)
        self.btn_office.toggled.connect(self._on_mode_changed)
        self.btn_pdf.toggled.connect(self.changed)
        self.btn_office.toggled.connect(self.changed)
        self.format_combo.currentIndexChanged.connect(self._on_format_changed)
        self.format_combo.currentIndexChanged.connect(self.changed)
        self.dpi_spin.valueChanged.connect(self.changed)
        self.range_edit.textChanged.connect(self.changed)
        self.xlsx_check.stateChanged.connect(self.changed)

        self._on_mode_changed()

    def _on_mode_changed(self):
        """操作模式切换"""
        is_pdf_mode = self.btn_pdf.isChecked()
        self.format_combo.clear()
        if is_pdf_mode:
            self.format_combo.addItems(["pdf", "docx", "xlsx", "pptx", "jpg", "png", "txt", "html"])
        else:
            self.format_combo.addItems(["docx", "doc","xlsx",  "xls", "pptx", "ppt"])
        self._on_format_changed()

    def _on_format_changed(self):
        """目标格式切换"""
        target = self.format_combo.currentText()
        is_pdf_mode = self.btn_pdf.isChecked()
        is_image = target in ("jpg", "png")
        is_xlsx = (target == "xlsx") or (not is_pdf_mode and target == "xlsx")
        is_office_convert = not is_pdf_mode
        need_range = is_pdf_mode and target != "pdf"
        self.dpi_label.setVisible(is_image)
        self.dpi_spin.setVisible(is_image)
        self.xlsx_check.setVisible(is_xlsx)
        self.range_widget.setVisible(need_range)


def build_panel() -> QWidget:
    """构建面板实例"""
    return ConvertPanel()


def collect_settings(panel: ConvertPanel) -> dict:
    """收集面板设置"""
    is_pdf_mode = panel.btn_pdf.isChecked()
    target = panel.format_combo.currentText()
    if is_pdf_mode:
        if target == "pdf":
            direction = "to_pdf" 
        else:
            direction = "from_pdf" 
    else:
        direction = "office_convert" 
    return {
        "direction": direction,
        "target": target,
        "dpi": panel.dpi_spin.value(),
        "range": panel.range_edit.text().strip(),
        "xlsx_merge_check_to_one_sheet": panel.xlsx_check.isChecked(),
    }


def prepare_preview(items, settings):
    """生成预览信息"""
    direction = settings.get("direction", "to_pdf")
    target = settings.get("target", "pdf")
    for it in items:
        dpi = settings.get("dpi", 72)
        pg = settings.get("range", "")
        xlsx_merge = settings.get("xlsx_merge_check_to_one_sheet", False)
        if direction == "to_pdf":
            hint = f"→ pdf"
        elif direction == "office_convert":
            hint = f"Office 互转 → {target}"
        else: 
            if target in ("jpg", "png"):
                hint = f"pdf → {target}，DPI={dpi}"
                if pg:
                    hint += f"，页码{pg}"
            elif target == "xlsx":
                hint = f"pdf → xlsx"
                if xlsx_merge:
                    hint += "，合并为单表"
                if pg:
                    hint += f"，页码{pg}"
            else:
                hint = f"pdf → {target}"
                if pg:
                    hint += f"，页码{pg}"
        it.preview_extra = {"A": hint}


def _replace_ext_keep_basename(name: str, new_ext: str) -> str:
    """替换文件扩展名，保留基础名称"""
    base, _ = os.path.splitext(name)
    if not new_ext.startswith("."):
        new_ext = "." + new_ext
    return base + new_ext


def _safe_ensure_dir(path: str):
    """确保目录存在"""
    os.makedirs(path, exist_ok=True)


def _is_image_ext(ext: str) -> bool:
    """判断是否为图片扩展名"""
    return ext.lower() in (".jpg", ".jpeg", ".png", ".bmp", ".webp", ".tif", ".tiff")


def _is_office_ext(ext: str) -> bool:
    """判断是否为 Office 文档扩展名"""
    return ext.lower() in (".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx")


def _which(cmd: str):
    """查找可执行文件"""
    from shutil import which
    return which(cmd)


def _office_to_pdf_windows(input_path: str, out_path: str):
    """使用 Microsoft Office COM 将 Office 文档转换为 PDF"""
    import pythoncom
    import comtypes.client
    input_path = urllib.parse.unquote(input_path)
    input_path = os.path.abspath(os.path.normpath(input_path))
    out_path = os.path.abspath(os.path.normpath(out_path))
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"文件不存在: {input_path}")
    need_uninit = False
    app = None
    doc = None
    try:
        try:
            pythoncom.CoGetInterfaceAndReleaseStream(pythoncom.IID_IDispatch, 0)
        except:
            pythoncom.CoInitialize()
            need_uninit = True

        ext = os.path.splitext(input_path)[1].lower()
        if ext in (".doc", ".docx"):
            app = comtypes.client.CreateObject("Word.Application")
            app.Visible = False
            app.DisplayAlerts = 0
            doc = app.Documents.Open(input_path, ReadOnly=True)
            doc.ExportAsFixedFormat(out_path, 17)
            doc.Close()
        elif ext in (".xls", ".xlsx"):
            app = comtypes.client.CreateObject("Excel.Application")
            app.Visible = False
            app.DisplayAlerts = False
            wb = app.Workbooks.Open(input_path)
            wb.ExportAsFixedFormat(0, out_path)
            wb.Close(False)
        elif ext in (".ppt", ".pptx"):
            app = comtypes.client.CreateObject("PowerPoint.Application")
            app.DisplayAlerts = False
            pres = app.Presentations.Open(input_path, WithWindow=False)
            pres.SaveAs(out_path, 32)
            pres.Close()
        else:
            raise RuntimeError(f"不支持的 Office 格式: {ext}")
    except Exception as e:
        raise RuntimeError(f"Microsoft Office 转换失败: {e}\n请检查文件是否可读，或尝试以管理员身份运行。")
    finally:
        if doc:
            try:
                doc.Close()
            except:
                pass
        if app:
            try:
                app.Quit()
            except:
                pass
        if need_uninit:
            pythoncom.CoUninitialize()


def _office_to_pdf_libreoffice(input_path: str, out_path: str):
    """使用 LibreOffice 将文档转换为 PDF"""
    input_path = urllib.parse.unquote(input_path)
    input_path = os.path.normpath(input_path)
    soffice = _which("soffice") or _which("libreoffice")
    if not soffice:
        raise RuntimeError("未检测到 LibreOffice，请安装 LibreOffice 或将 soffice 加入 PATH")
    out_dir = os.path.dirname(out_path) or os.getcwd()
    _safe_ensure_dir(out_dir)
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, input_path],
                   check=True, capture_output=True, text=True, encoding='utf-8', errors='ignore')
    base = os.path.splitext(os.path.basename(input_path))[0]
    produced = os.path.join(out_dir, base + ".pdf")
    if os.path.abspath(produced) != os.path.abspath(out_path):
        if os.path.exists(out_path):
            os.remove(out_path)
        os.replace(produced, out_path)


def _image_to_pdf(input_path: str, out_path: str):
    """将图片转换为 PDF"""
    if Image is None:
        raise RuntimeError("缺少 Pillow 库")
    with Image.open(input_path) as im:
        if im.mode in ("RGBA", "LA"):
            im = im.convert("RGB")
        im.save(out_path, "PDF")


def _pdf_to_text(input_path: str) -> str:
    """提取 PDF 文本内容"""
    if fitz is None:
        raise RuntimeError("缺少 PyMuPDF")
    doc = fitz.open(input_path)
    parts = []
    for i, page in enumerate(doc):
        parts.append(f"\n\n---- Page {i+1} ----\n")
        parts.append(page.get_text("text"))
    doc.close()
    return "".join(parts)


def _pdf_to_html(input_path: str) -> str:
    """将 PDF 转换为 XHTML 字符串"""
    if fitz is None:
        raise RuntimeError("缺少 PyMuPDF")
    doc = fitz.open(input_path)
    body = [page.get_text("xhtml") for page in doc]
    doc.close()
    return "<!DOCTYPE html><html><head><meta charset='utf-8'><title>PDF Export</title></head><body>" + \
           "\n".join(body) + "</body></html>"


def _pdf_to_images_like_reference(input_path, out_dir, base_name, fmt, dpi, page_expr, custom_names=None):
    """将 PDF 页面导出为图片"""
    if fitz is None:
        raise RuntimeError("缺少 PyMuPDF")
    base_name = base_name or os.path.splitext(os.path.basename(input_path))[0]
    doc = fitz.open(input_path)
    written = []
    try:
        if doc.needs_pass:
            raise RuntimeError("PDF 受密码保护")
        base_zoom = max(72, min(600, dpi)) / 72.0
        indices = parse_page_range(page_expr, len(doc))
        for i, idx in enumerate(indices):
            page = doc.load_page(idx)
            zoom = base_zoom
            if max(page.rect.width * zoom, page.rect.height * zoom) > 8000:
                zoom = 8000 / max(page.rect.width, page.rect.height)
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
            if custom_names and i < len(custom_names):
                name = custom_names[i]
            else:
                name = f"{base_name}_page_{i+1}"
            out_file = os.path.join(out_dir, f"{name}.{fmt.lower()}")
            pix.save(out_file)
            written.append(out_file)
        return written
    finally:
        doc.close()


def _pdf_to_excel(input_path: str, out_path: str, page_expr: str = "", merge_to_one_sheet: bool = False):
    """将 PDF 表格导出为 Excel"""
    try:
        import pdfplumber
        import pandas as pd
    except ImportError:
        raise RuntimeError("缺少 pdfplumber 或 pandas，请执行: pip install pdfplumber pandas openpyxl")
    with pdfplumber.open(input_path) as pdf:
        pages = parse_page_range(page_expr, len(pdf.pages))
        writer = pd.ExcelWriter(out_path, engine="openpyxl")
        try:
            if merge_to_one_sheet:
                rows = []
                for idx in pages:
                    tables = pdf.pages[idx].extract_tables()
                    for t in tables or []:
                        rows.extend(t + [[""]])
                pd.DataFrame(rows).to_excel(writer, sheet_name="Tables", index=False, header=False)
            else:
                for idx in pages:
                    page = pdf.pages[idx]
                    tables = page.extract_tables()
                    if tables:
                        for ti, t in enumerate(tables, 1):
                            pd.DataFrame(t).to_excel(writer, sheet_name=f"P{idx+1}_T{ti}", index=False, header=False)
                    else:
                        txt = page.extract_text() or ""
                        pd.DataFrame([[line] for line in txt.splitlines()]).to_excel(writer, sheet_name=f"P{idx+1}_Text", index=False, header=False)
        finally:
            writer.close()


def _pdf_to_pptx(input_path: str, out_path: str, dpi: int = 150, page_expr: str = ""):
    """将 PDF 页面转换为 PPTX"""
    if fitz is None:
        raise RuntimeError("缺少 PyMuPDF")
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise RuntimeError("缺少 python-pptx，请执行: pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    doc = fitz.open(input_path)
    try:
        indices = parse_page_range(page_expr, len(doc))
        for idx in indices:
            page = doc[idx]
            zoom = max(72, min(600, dpi)) / 72.0
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
            width_in = pix.width / (dpi if dpi>0 else 150)
            height_in = pix.height / (dpi if dpi>0 else 150)
            prs.slide_width = Inches(width_in)
            prs.slide_height = Inches(height_in)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(io.BytesIO(pix.tobytes("png")), 0, 0, width=Inches(width_in), height=Inches(height_in))
    finally:
        doc.close()
    prs.save(out_path)


def _office_convert_windows(input_path: str, out_path: str, target_ext: str):
    """使用 Microsoft Office 进行文档格式转换"""
    import pythoncom
    import comtypes.client
    input_path = os.path.abspath(os.path.normpath(input_path))
    out_path = os.path.abspath(os.path.normpath(out_path))
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"文件不存在: {input_path}")
    need_uninit = False
    app = None
    doc = None
    try:
        try:
            pythoncom.CoGetInterfaceAndReleaseStream(pythoncom.IID_IDispatch, 0)
        except:
            pythoncom.CoInitialize()
            need_uninit = True
        ext = os.path.splitext(input_path)[1].lower()
        target_ext_lower = target_ext.lower()
        target_ext_clean = target_ext_lower.lstrip(".")
        if ext in (".doc", ".docx"):
            app = comtypes.client.CreateObject("Word.Application")
            app.Visible = False
            app.DisplayAlerts = 0
            doc = app.Documents.Open(input_path, ReadOnly=True)
            if target_ext_clean in ("docx", "doc"):
                doc.SaveAs(out_path, wdFormatDocx if target_ext_clean == "docx" else wdFormatDoc)
            elif target_ext_clean == "pdf":
                doc.ExportAsFixedFormat(out_path, 17)
            else:
                raise RuntimeError(f"Word 不支持转换为 {target_ext}")
            doc.Close()
        elif ext in (".xls", ".xlsx"):
            app = comtypes.client.CreateObject("Excel.Application")
            app.Visible = False
            app.DisplayAlerts = False
            wb = app.Workbooks.Open(input_path)
            if target_ext_clean in ("xlsx", "xls"):
                wb.SaveAs(out_path, xlOpenXMLWorkbook if target_ext_clean == "xlsx" else xlExcel8)
            elif target_ext_clean == "pdf":
                wb.ExportAsFixedFormat(0, out_path)
            else:
                raise RuntimeError(f"Excel 不支持转换为 {target_ext}")
            wb.Close(False)
        elif ext in (".ppt", ".pptx"):
            app = comtypes.client.CreateObject("PowerPoint.Application")
            app.DisplayAlerts = False
            pres = app.Presentations.Open(input_path, WithWindow=False)
            if target_ext_clean in ("pptx", "ppt"):
                pres.SaveAs(out_path, ppSaveAsPPTX if target_ext_clean == "pptx" else ppSaveAsPPT)
            elif target_ext_clean == "pdf":
                pres.SaveAs(out_path, 32)
            else:
                raise RuntimeError(f"PowerPoint 不支持转换为 {target_ext}")
            pres.Close()
        else:
            raise RuntimeError(f"不支持的文件格式: {ext}")
    except Exception as e:
        raise RuntimeError(f"Office 转换失败: {e}")
    finally:
        if doc:
            try:
                doc.Close()
            except:
                pass
        if app:
            try:
                app.Quit()
            except:
                pass
        if need_uninit:
            pythoncom.CoUninitialize()


def run_task(file_item, settings: dict, custom_names=None):
    """执行单个文件转换任务"""
    direction = settings.get("direction", "to_pdf")
    target = settings.get("target", "PDF")
    dpi = int(settings.get("dpi", 72))
    page_expr = settings.get("range", "")
    xlsx_merge_check = bool(settings.get("xlsx_merge_check_to_one_sheet", False))
    src = file_item.input_path
    in_ext = os.path.splitext(src)[1].lower()
    out_dir = file_item.output_dir or os.path.dirname(src)
    _safe_ensure_dir(out_dir)
    out_path = os.path.join(out_dir, file_item.output_name)
    if direction == "to_pdf":
        if in_ext == ".pdf":
            if os.path.abspath(src) != os.path.abspath(out_path):
                shutil.copy2(src, out_path)
        elif _is_image_ext(in_ext):
            _image_to_pdf(src, out_path)
        elif _is_office_ext(in_ext):
            if sys.platform.startswith("win") and comtypes is not None:
                try:
                    _office_to_pdf_windows(src, out_path)
                except Exception as e:
                    print(f"Office COM 转换失败，尝试 LibreOffice: {e}")
                    _office_to_pdf_libreoffice(src, out_path)
            else:
                _office_to_pdf_libreoffice(src, out_path)
        else:
            raise RuntimeError(f"不支持的文件格式: {in_ext}，无法转换为 PDF")
        file_item.output_name = os.path.basename(out_path)
    elif direction == "office_convert":
        if not _is_office_ext(in_ext):
            raise RuntimeError(f"不支持的文件格式: {in_ext}，仅支持 Word/Excel/PPT 文档")
        if sys.platform.startswith("win") and comtypes is not None:
            output_ext = target.lower()
            if in_ext == f".{output_ext}":
                shutil.copy2(src, out_path)
            else:
                _office_convert_windows(src, out_path, output_ext)
        else:
            raise RuntimeError("Office 文档互转需要 Windows + Microsoft Office")
    else: 
        if in_ext != ".pdf":
            shutil.copy2(src, out_path)
            file_item.output_name = os.path.basename(out_path)
        else:
            target_lower = target.lower()
            if target_lower == "txt":
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write(_pdf_to_text(src))
            elif target_lower == "html":
                with open(out_path, "w", encoding="utf-8") as f:
                    f.write(_pdf_to_html(src))
            elif target_lower == "docx":
                if Converter is None:
                    raise RuntimeError("缺少 pdf2docx 库")
                cv = Converter(src)
                try:
                    cv.convert(out_path)
                finally:
                    cv.close()
            elif target_lower in ("jpg", "png"):
                base = os.path.splitext(file_item.output_name)[0] if file_item.output_name else os.path.splitext(os.path.basename(src))[0]
                written = _pdf_to_images_like_reference(src, out_dir, base, target_lower, dpi, page_expr, custom_names=settings.get("custom_names"))
                file_item.output_paths = written
                file_item.output_name = os.path.basename(written[0]) if written else ""
            elif target_lower == "xlsx":
                _pdf_to_excel(src, out_path, page_expr, xlsx_merge_check)
            elif target_lower == "pptx":
                _pdf_to_pptx(src, out_path, dpi, page_expr)
            else:
                raise RuntimeError(f"不支持的目标格式: {target}")
            if target_lower not in ("jpg", "png"):
                file_item.output_name = os.path.basename(out_path)
    file_item.status = "完成"